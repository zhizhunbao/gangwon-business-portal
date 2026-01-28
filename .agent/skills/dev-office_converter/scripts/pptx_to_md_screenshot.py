import os
import sys
import argparse
from pathlib import Path
import win32com.client
from pptx import Presentation

def pptx_to_md_with_screenshots(pptx_path: Path, output_md: Path = None, dpi: int = 200):
    """
    Directly convert PPTX to Markdown with screenshots for each slide.
    Uses PowerPoint COM interface for high-quality rendering.
    """
    if not output_md:
        output_md = pptx_path.with_suffix(".md")
    
    # Create directory for screenshots
    assets_dir = output_md.parent / f"{output_md.stem}_assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"🚀 Starting PowerPoint for: {pptx_path.name}")
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    except Exception as e:
        print(f"❌ Error: Could not start PowerPoint. Make sure it's installed. {e}")
        return

    try:
        abs_pptx = str(pptx_path.absolute())
        presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        
        md_content = []
        md_content.append(f"# {pptx_path.stem}\n")
        md_content.append(f"**Source:** `{pptx_path.name}`\n")
        md_content.append("---\n")
        
        # Load via python-pptx for text extraction (sometimes cleaner)
        prs_text = Presentation(pptx_path)
        
        for i, slide in enumerate(presentation.Slides, 1):
            print(f"  📸 Capturing Slide {i}/{presentation.Slides.Count}...")
            
            # 1. Export Screenshot
            image_name = f"slide_{i:03d}.png"
            image_path = assets_dir / image_name
            # Export takes (FileName, FilterName, ScaleWidth, ScaleHeight)
            # Scaling for DPI (96 is default)
            scale = dpi / 96.0
            width = presentation.PageSetup.SlideWidth * scale
            height = presentation.PageSetup.SlideHeight * scale
            slide.Export(str(image_path.absolute()), "PNG", width, height)
            
            # 2. Extract Text (via COM or python-pptx)
            # We'll use python-pptx for text as it's often more reliable for structural text
            slide_text_obj = prs_text.slides[i-1]
            title = ""
            if slide_text_obj.shapes.title:
                title = slide_text_obj.shapes.title.text.strip()
            
            md_content.append(f"## Slide {i}: {title}\n")
            md_content.append(f"![Slide {i}]({assets_dir.name}/{image_name})\n")
            
            # Extract bullet points/text
            text_runs = []
            for shape in slide_text_obj.shapes:
                if shape == slide_text_obj.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if '\n' in text:
                        lines = [f"- {line.strip()}" for line in text.split('\n') if line.strip()]
                        text_runs.extend(lines)
                    else:
                        text_runs.append(text)
            
            if text_runs:
                md_content.append("### 📝 Text Content\n")
                md_content.append("\n".join(text_runs))
                md_content.append("\n")
            
            md_content.append("---\n")
            
        presentation.Close()
        
        # Save Markdown
        output_md.write_text("\n".join(md_content), encoding='utf-8')
        print(f"\n✨ Successfully converted to: {output_md}")
        print(f"📦 Screenshots saved in: {assets_dir}/")
        
    except Exception as e:
        print(f"❌ Critical Error: {e}")
        import traceback
        traceback.print_exc()
    finally:
        powerpoint.Quit()

def main():
    parser = argparse.ArgumentParser(description="Convert PPTX directly to Markdown with screenshots")
    parser.add_argument("input", help="Input PPTX file path")
    parser.add_argument("-o", "--output", help="Output MD file path")
    parser.add_argument("--dpi", type=int, default=200, help="Output image DPI (default: 200)")
    
    args = parser.parse_args()
    
    input_pptx = Path(args.input)
    if not input_pptx.exists():
        print(f"❌ Error: File {input_pptx} not found.")
        sys.exit(1)
        
    output_md = Path(args.output) if args.output else None
    pptx_to_md_with_screenshots(input_pptx, output_md, args.dpi)

if __name__ == "__main__":
    main()
